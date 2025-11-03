# ---------------------------------------------------------------
# 📄 Docify RAG + Extended Evaluation Metrics + Combined Heatmap
# ---------------------------------------------------------------

import os
import tempfile
import io
from typing import List, Tuple
import streamlit as st
from PIL import Image
import fitz  # PyMuPDF
from pdf2image import convert_from_path
import pytesseract
import docx
from odf import text as odf_text
from odf.opendocument import load as odf_load
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize
nltk.download('punkt', quiet=True)

from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import torch

# Metrics and visualization packages
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from rouge_score import rouge_scorer
from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
from nltk.translate.meteor_score import meteor_score

# Try to use optional BERTScore if available
try:
    from bert_score import score as bert_score_fn
    BERTSCORE_AVAILABLE = True
except Exception:
    BERTSCORE_AVAILABLE = False

# -----------------------------
# Configuration / Models
# -----------------------------
EMBEDDING_MODEL_NAME = 'all-MiniLM-L6-v2'
GEN_MODEL_NAME = 'google/flan-t5-small'
DEVICE = 'cuda' if torch.cuda.is_available() else 'cpu'

@st.cache_resource
def load_models():
    embedder = SentenceTransformer(EMBEDDING_MODEL_NAME, device=DEVICE)
    tokenizer = AutoTokenizer.from_pretrained(GEN_MODEL_NAME)
    gen_model = AutoModelForSeq2SeqLM.from_pretrained(GEN_MODEL_NAME).to(DEVICE)
    return embedder, tokenizer, gen_model

embedder, tokenizer, gen_model = load_models()

# -----------------------------
# File Extraction Functions
# -----------------------------
def extract_text_from_pdf(path: str) -> str:
    try:
        doc = fitz.open(path)
        text = "\n".join(page.get_text() for page in doc)
        if text.strip():
            return text
        images = convert_from_path(path)
        return "\n".join(pytesseract.image_to_string(img) for img in images)
    except Exception:
        return ""

def extract_text_from_docx(path: str) -> str:
    try:
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""

def extract_text_from_odt(path: str) -> str:
    try:
        doc = odf_load(path)
        paras = doc.getElementsByType(odf_text.P)
        return "\n".join(p.firstChild.data if p.firstChild else '' for p in paras)
    except Exception:
        return ""

def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""

def extract_text_from_image(path: str) -> str:
    try:
        return pytesseract.image_to_string(Image.open(path))
    except Exception:
        return ""

# -----------------------------
# Text Preprocessing
# -----------------------------
def preprocess_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace('\r', '\n')
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())

def chunk_text(text: str, max_tokens: int = 200) -> List[str]:
    sents = sent_tokenize(text)
    chunks, cur, cur_len = [], [], 0
    for s in sents:
        l = len(s.split())
        if cur_len + l > max_tokens and cur:
            chunks.append(" ".join(cur))
            cur, cur_len = [s], l
        else:
            cur.append(s)
            cur_len += l
    if cur:
        chunks.append(" ".join(cur))
    return chunks

# -----------------------------
# RAG Index + Retrieval
# -----------------------------
def build_faiss_index(docs: List[str]):
    embeddings = embedder.encode(docs, convert_to_numpy=True, show_progress_bar=False)
    dim = embeddings.shape[1]
    faiss.normalize_L2(embeddings)
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)
    return index, embeddings

def retrieve(index, embeddings, docs: List[str], query: str, top_k: int = 5):
    q_emb = embedder.encode([query], convert_to_numpy=True)
    faiss.normalize_L2(q_emb)
    D, I = index.search(q_emb, top_k)
    return [(int(idx), float(score), docs[idx]) for score, idx in zip(D[0], I[0]) if idx < len(docs)]

# -----------------------------
# RAG Generation
# -----------------------------
def rag_generate(query: str, docs: List[str], index, embeddings, top_k: int = 3):
    hits = retrieve(index, embeddings, docs, query, top_k)
    context = "\n---\n".join(h[2] for h in hits)
    prompt = f"Answer the query using the context below. If unknown, say 'Information not found.'\n\nContext:\n{context}\n\nQuery: {query}\nAnswer:"
    inputs = tokenizer(prompt, return_tensors="pt", truncation=True, max_length=1024).to(DEVICE)
    out = gen_model.generate(**inputs, max_length=256, num_beams=4)
    answer = tokenizer.decode(out[0], skip_special_tokens=True)
    return answer, hits

# -----------------------------
# Extended Metric Computation
# -----------------------------
smoothie = SmoothingFunction().method4
rouge = rouge_scorer.RougeScorer(["rouge1", "rouge2", "rougeL"], use_stemmer=True)

def safe_bleu(reference, hypothesis):
    try:
        return sentence_bleu([reference.split()], hypothesis.split(), smoothing_function=smoothie)
    except Exception:
        return 0.0

def safe_meteor(reference, hypothesis):
    try:
        return meteor_score([reference], hypothesis)
    except Exception:
        return 0.0

def compute_bertscore_list(refs, hyps, lang="en"):
    if not BERTSCORE_AVAILABLE:
        return [np.nan] * len(refs)
    P, R, F = bert_score_fn(hyps, refs, lang=lang, verbose=False, rescale_with_baseline=True)
    return [float(f) for f in F]

def overlap_ratio(text_a, text_b):
    a_tokens = set(text_a.lower().split())
    b_tokens = set(text_b.lower().split())
    if not a_tokens:
        return 0.0
    return len(a_tokens & b_tokens) / len(a_tokens)

def composite_score(row):
    cos = row.get("cosine", 0.0)
    r1 = row.get("rouge1", 0.0)
    bs = row.get("bertscore", np.nan)
    if not np.isnan(bs):
        return 0.4 * cos + 0.3 * r1 + 0.3 * bs
    else:
        b = row.get("bleu", 0.0)
        return 0.5 * cos + 0.3 * r1 + 0.2 * b

# -----------------------------
# Per-chunk metrics and Heatmap
# -----------------------------
def compute_extended_metrics_for_chunks(chunks: List[str], answer: str, chunk_embs: np.ndarray):
    n = len(chunks)
    metrics = {m: [] for m in ["rouge1","rouge2","rougeL","bleu","meteor","cosine","overlap"]}
    ans = answer or ""
    for i, chunk in enumerate(chunks):
        r = rouge.score(chunk, ans)
        metrics["rouge1"].append(r["rouge1"].fmeasure)
        metrics["rouge2"].append(r["rouge2"].fmeasure)
        metrics["rougeL"].append(r["rougeL"].fmeasure)
        metrics["bleu"].append(safe_bleu(chunk, ans))
        metrics["meteor"].append(safe_meteor(chunk, ans))
        ans_emb = embedder.encode([ans], convert_to_numpy=True)
        cos = float(cosine_similarity(chunk_embs[i:i+1], ans_emb)[0,0])
        metrics["cosine"].append(cos)
        metrics["overlap"].append(overlap_ratio(ans, chunk))
    metrics["bertscore"] = compute_bertscore_list(chunks, [ans]*n, lang="en") if BERTSCORE_AVAILABLE else [np.nan]*n
    cos_arr = np.array(metrics["cosine"])
    topk_order = np.argsort(-cos_arr)
    df = pd.DataFrame(metrics)
    df.index = [f"Chunk {i+1}" for i in range(n)]
    ranks = np.empty(n, dtype=int)
    ranks[topk_order] = np.arange(1, n+1)
    df["relevance_rank"] = ranks
    df["relevance_norm"] = 1.0 - (ranks - 1) / max(1, n-1)
    df["composite"] = df.apply(composite_score, axis=1)
    return df

def build_combined_heatmap(df_chunks: pd.DataFrame, title: str):
    doc_series = df_chunks.mean(axis=0, numeric_only=True)
    doc_series.name = "Document"
    combined = df_chunks.copy()
    combined.loc["Document"] = doc_series
    numeric_cols = combined.columns.tolist()
    heat_df = combined[numeric_cols].T
    heat_norm = heat_df.copy()
    for r in heat_df.index:
        row = heat_df.loc[r].astype(float)
        mn, mx = row.min(), row.max()
        heat_norm.loc[r] = (row - mn) / (mx - mn) if mx > mn else 0.5
    ordered = ["composite", "cosine", "rouge1", "rouge2", "rougeL", "bleu", "meteor", "bertscore", "overlap", "relevance_norm", "relevance_rank"]
    ordered = [m for m in ordered if m in heat_norm.index] + [m for m in heat_norm.index if m not in ordered]
    heat_norm = heat_norm.loc[ordered]
    plt.figure(figsize=(max(10, 1.2 * heat_norm.shape[1]), max(6, 0.5 * heat_norm.shape[0])))
    sns.heatmap(heat_norm, annot=True, fmt=".2f", cmap="viridis", cbar_kws={"label": "Normalized metric (0-1)"})
    plt.title(title)
    plt.xlabel("Chunks → Document")
    plt.ylabel("Metrics")
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    st.pyplot(plt.gcf())
    return combined, heat_norm

# -----------------------------
# Streamlit Interface
# -----------------------------
st.set_page_config(page_title="📄 Docify RAG with Extended Metrics", layout="wide")
st.title("📄 Docify — Document RAG + Advanced Metrics + Heatmap")

uploaded = st.file_uploader("Upload document", type=["pdf", "docx", "odt", "pptx", "txt", "jpg", "png"])
max_chunk = st.slider("Max chunk tokens", 50, 800, 200, 10)

if uploaded:
    suffix = uploaded.name.split(".")[-1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix="."+suffix) as tmp:
        tmp.write(uploaded.read())
        path = tmp.name

    # Text extraction dispatch
    if suffix == "pdf":
        raw_text = extract_text_from_pdf(path)
    elif suffix == "docx":
        raw_text = extract_text_from_docx(path)
    elif suffix == "odt":
        raw_text = extract_text_from_odt(path)
    elif suffix == "pptx":
        raw_text = extract_text_from_pptx(path)
    elif suffix == "txt":
        raw_text = extract_text_from_txt(path)
    elif suffix in ["jpg", "jpeg", "png"]:
        raw_text = extract_text_from_image(path)
    else:
        raw_text = ""

    clean_text = preprocess_text(raw_text)
    chunks = chunk_text(clean_text, max_chunk)
    index, emb = build_faiss_index(chunks)

    st.subheader("Extracted & Chunked Text")
    st.text_area("Processed Text", clean_text[:5000], height=250)
    st.write(f"📚 Total chunks: {len(chunks)}")

    st.markdown("---")
    query = st.text_input("🔍 Enter your question or query:")
    if query:
        with st.spinner("Generating answer using RAG..."):
            ans, hits = rag_generate(query, chunks, index, emb)
        st.subheader("Generated Answer")
        st.write(ans)

        if hits:
            try:
                chunk_embs = emb if emb.shape[0] == len(chunks) else embedder.encode(chunks, convert_to_numpy=True)
            except Exception:
                chunk_embs = embedder.encode(chunks, convert_to_numpy=True)

            with st.spinner("Computing extended metrics and combined heatmap..."):
                df_chunks = compute_extended_metrics_for_chunks(chunks, ans, chunk_embs)
                combined_df, norm_df = build_combined_heatmap(df_chunks, "All Metrics (Chunks → Document)")

            st.subheader("📊 Document-level aggregated metrics (combined)")
            doc_row = combined_df.loc["Document"].to_dict()
            st.json({k: (None if pd.isna(v) else float(np.round(v, 4))) for k, v in doc_row.items()})

            st.subheader("📋 Raw per-chunk metric table (hidden by default)")
            with st.expander("Show per-chunk metrics table"):
                st.dataframe(df_chunks.style.format("{:.4f}"))